import warnings
import sys
import os
import io

# Disable all warnings (as suggested by GLiNER owner)
warnings.filterwarnings("ignore")

# Create a more aggressive stdout suppressor
class SuppressStdoutStderr:
    def __init__(self):
        self.null_fds = [os.open(os.devnull, os.O_RDWR) for x in range(2)]
        self.save_fds = [os.dup(1), os.dup(2)]

    def __enter__(self):
        os.dup2(self.null_fds[0], 1)
        os.dup2(self.null_fds[1], 2)

    def __exit__(self, *_):
        os.dup2(self.save_fds[0], 1)
        os.dup2(self.save_fds[1], 2)
        for fd in self.null_fds + self.save_fds:
            os.close(fd)

warnings.filterwarnings("ignore", category=UserWarning)  # General UserWarnings
warnings.filterwarnings("ignore", module="gliner")  # All GLiNER warnings
warnings.filterwarnings("ignore", module="transformers")  # All transformers warnings
warnings.filterwarnings("ignore", message=".*truncate to max_length.*", category=UserWarning)  # Specific truncation warning
warnings.filterwarnings("ignore", message=".*no maximum length is provided.*", category=UserWarning)  # Another form of the warning

os.environ["ONNXRUNTIME_DISABLE_VERSION_WARNING"] = "1"
os.environ["TQDM_DISABLE"] = "1"
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["TRANSFORMERS_NO_ADVISORY_WARNINGS"] = "1"

import sqlite3
import hashlib
from datetime import datetime, timezone
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from gliner import GLiNER
from termcolor import colored
from dotenv import load_dotenv
import argparse
import logging

# Add LOG_LEVEL definition
LOG_LEVEL = os.getenv("LOG_LEVEL", "INFO")
LOG_FILE = os.path.join(os.environ.get("PROGRAMDATA", ""), "PII Scanner", "pii_scanner.log")

# Simplify logging setup
handler = logging.StreamHandler(sys.stdout)
handler.setLevel(logging.INFO)

# Configure logging with handler
logging.basicConfig(
    level=getattr(logging, LOG_LEVEL.upper()),
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[handler]
)

logger = logging.getLogger(__name__)

# Disable tqdm progress bars


def get_env_file_path():
    """Get the path to the .env file based on whether running as exe or script"""
    if getattr(sys, 'frozen', False):
        # Running as exe
        return os.path.join(os.path.dirname(sys.executable), '.env')
    else:
        # Running as script
        return os.path.join(os.path.dirname(os.path.abspath(__file__)), '.env')

# Load environment variables from .env file
env_path = get_env_file_path()
if os.path.exists(env_path):
    load_dotenv(env_path)
    logger.info(f"Loaded configuration from {env_path}")
else:
    logger.warning(f"No .env file found at {env_path}, using defaults")

def get_application_path():
    """Get the base path for the application, works in both script and exe"""
    if getattr(sys, 'frozen', False):
        # Running as compiled exe
        return os.path.dirname(sys.executable)
    else:
        # Running as script
        return os.path.dirname(os.path.abspath(__file__))

# Configuration - Pull from environment variables or use defaults
DB_FILE = os.getenv("DB_FILE", os.path.join(os.environ.get("PROGRAMDATA", ""),
                                           "PII Scanner",
                                           "pii_scan_history.db"))
ALLOWED_FILE_TYPES = {".doc", ".docx", ".xlsx", ".pptx", ".txt"}
MAX_CHUNK_LENGTH = int(os.getenv("MAX_CHUNK_LENGTH", "384"))  # Changed default to 384
LITE_SCAN_LIMIT = 1024 * 1024  # 1MB limit for lite scan
PII_MODEL_NAME = os.getenv("PII_MODEL_NAME", "urchade/gliner_multi_pii-v1")

# Default PII labels if not specified in .env
DEFAULT_PII_LABELS = [
    "person", "organization", "phone number", "address", "passport number",
    "email", "credit card number", "social security number"
]

# Load PII labels from environment
PII_LABELS = os.getenv("PII_LABELS", ",".join(DEFAULT_PII_LABELS)).split(",")
PII_LABELS_FULL = os.getenv("PII_LABELS_FULL", ",".join(PII_LABELS)).split(",")

# Print configured labels
print(colored("\nConfigured PII labels:", "cyan"))
print(colored("Basic labels:", "yellow"))
for label in PII_LABELS:
    print(colored(f"  - {label}", "yellow"))

print(colored("\nFull label set available:", "green"))
print(colored(f"  Total labels: {len(PII_LABELS_FULL)}", "green"))

# Add these constants after the existing config section
EXIT_SUCCESS = 0
EXIT_PII_FOUND = 1
EXIT_FILE_NOT_FOUND = 2
EXIT_UNSUPPORTED_FILE = 3
EXIT_DB_ERROR = 4
EXIT_TOKENIZER_INIT_ERROR = 5
EXIT_GLINER_INIT_ERROR = 6
EXIT_NLTK_INIT_ERROR = 7
EXIT_CHECKSUM_ERROR = 8
EXIT_TEXT_EXTRACTION_ERROR = 9
EXIT_TEXT_CHUNKING_ERROR = 10
EXIT_PII_DETECTION_ERROR = 11
EXIT_INVALID_SCAN_TYPE = 12
EXIT_MISSING_PATH = 13
EXIT_GENERAL_ERROR = 99

# Ensure directories exist
os.makedirs(os.path.dirname(LOG_FILE), exist_ok=True)
os.makedirs(os.path.dirname(DB_FILE), exist_ok=True)

def init_nltk():
    """Initialize NLTK resources."""
    logger.info("NLTK initialization skipped - using GLiNER tokenizer")
    return

# Initialize GLiNER model with proper configuration
try:
    model_config = {
        'max_length': MAX_CHUNK_LENGTH,
        'padding': 'max_length',
        'truncation': True,
        'add_prefix_space': True
    }
    
    gliner_model = GLiNER.from_pretrained(
        PII_MODEL_NAME,
        **model_config
    )
    
    logger.info(f"GLiNER model '{PII_MODEL_NAME}' initialized successfully.")
except Exception as e:
    logger.error(f"Error initializing GLiNER model: {e}")
    gliner_model = None

def init_db():
    """Initialize the SQLite database."""
    conn = None
    try:
        conn = sqlite3.connect(DB_FILE)
        cursor = conn.cursor()
        
        # Create the scan_history table
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS scan_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                file_path TEXT NOT NULL,
                scan_time TEXT NOT NULL,
                file_size INTEGER,
                file_modified TEXT,
                file_checksum TEXT,
                scan_type TEXT CHECK(scan_type IN ('lite', 'full')),
                pii_entities TEXT,
                UNIQUE(file_path, scan_type)
            )
        """)
        conn.commit()
        logger.debug(f"Database initialized successfully at: {DB_FILE}")  # Changed to debug level
        
    except sqlite3.Error as e:
        logger.warning(f"Database initialization error: {e}")  # Changed to warning level
        if conn:
            conn.rollback()
        sys.exit(EXIT_DB_ERROR)
    finally:
        if conn:
            conn.close()

def calculate_checksum(file_path, scan_type):
    """Calculate SHA256 checksum of a file."""
    hasher = hashlib.sha256()
    try:
        with open(file_path, "rb") as f:
            if scan_type == "lite":
                chunk = f.read(LITE_SCAN_LIMIT)
                hasher.update(chunk)
            else:
                while True:
                    chunk = f.read(4096)
                    if not chunk:
                        break
                    hasher.update(chunk)
        return hasher.hexdigest()
    except Exception as e:
        logger.error(f"Error calculating checksum for {file_path}: {e}")
        return None

def is_file_scanned(file_path, checksum, scan_type):
    """Check if file has been scanned and return PII entities."""
    conn = None
    try:
        conn = sqlite3.connect(DB_FILE)
        cursor = conn.cursor()
        cursor.execute("""
            SELECT pii_entities FROM scan_history
            WHERE file_path = ? AND file_checksum = ? AND scan_type = ?
        """, (file_path, checksum, scan_type))
        result = cursor.fetchone()
        if result:
            pii_entities_str = result[0]
            if pii_entities_str:
                return True, eval(pii_entities_str)  # Returns True and the PII entities
            else:
                return True, []  # Already scanned, no PII found
        else:
            return False, None  # Not scanned yet
    except sqlite3.Error as e:
        logger.error(f"Database error checking if file is scanned: {e}")
        return False, None  # Treat as not scanned
    finally:
        if conn:
            conn.close()

def save_scan_result(file_path, pii_entities, file_size, file_modified, file_checksum, scan_type):
    """Save scan results to database with improved error handling."""
    conn = None
    try:
        conn = sqlite3.connect(DB_FILE)
        cursor = conn.cursor()
        now = datetime.now(timezone.utc).isoformat()
        
        # Validate inputs
        if not isinstance(file_size, int):
            file_size = int(file_size)
        
        if not scan_type in ['lite', 'full']:
            raise ValueError(f"Invalid scan_type: {scan_type}")
            
        if not file_checksum:
            raise ValueError("file_checksum cannot be None")
            
        # Print debug info
        logger.info("Saving scan result:")
        logger.info(f"  Path: {file_path}")
        logger.info(f"  Checksum: {file_checksum}")
        logger.info(f"  Scan type: {scan_type}")
            
        # Insert or update the scan result
        cursor.execute("""
            INSERT OR REPLACE INTO scan_history 
            (file_path, scan_time, file_size, file_modified, file_checksum, scan_type, pii_entities)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        """, (file_path, now, file_size, file_modified, file_checksum, scan_type, str(pii_entities)))
        
        conn.commit()
        logger.info(f"Scan result saved for: {file_path}")
        
    except sqlite3.Error as e:
        logger.error(f"Database error while saving scan result: {e}")
        logger.error(f"Error details: {str(e)}")
        if conn:
            conn.rollback()
        raise
    except ValueError as e:
        logger.error(f"Invalid data error: {e}")
        raise
    finally:
        if conn:
            conn.close()

def extract_text_from_file(file_path, scan_type):
    """Extract text from various file types."""
    try:
        if file_path.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read(LITE_SCAN_LIMIT) if scan_type == "lite" else f.read()
                
        elif file_path.endswith((".doc", ".docx")):
            document = Document(file_path)
            text = ""
            for paragraph in document.paragraphs:
                text += paragraph.text + "\n"
                if scan_type == "lite" and len(text.encode('utf-8')) > LITE_SCAN_LIMIT:
                    return text[:LITE_SCAN_LIMIT].decode('utf-8', 'ignore')
            return text
            
        elif file_path.endswith(".xlsx"):
            # Suppress openpyxl warnings during workbook load
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                workbook = load_workbook(file_path)
                text = ""
                for sheet in workbook:
                    for row in sheet.iter_rows():
                        text += " ".join([str(cell.value) if cell.value is not None else "" for cell in row]) + "\n"
                        if scan_type == "lite" and len(text.encode('utf-8')) > LITE_SCAN_LIMIT:
                            return text[:LITE_SCAN_LIMIT].decode('utf-8', 'ignore')
                return text
            
        elif file_path.endswith(".pptx"):
            presentation = Presentation(file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    if scan_type == "lite" and len(text.encode('utf-8')) > LITE_SCAN_LIMIT:
                        return text[:LITE_SCAN_LIMIT].decode('utf-8', 'ignore')
            return text
            
        return None
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return None

def chunk_text(text, max_length=MAX_CHUNK_LENGTH):
    """Split text into chunks using GLiNER's words_splitter."""
    try:
        if not text:
            return []
        
        # Get all tokens at once, extract only the token text from tuples
        tokens = [t[0] for t in gliner_model.data_processor.words_splitter(text)]
        
        # Initialize variables for chunking
        chunks = []
        current_chunk = []
        current_length = 0
        
        # Process tokens into chunks
        for token in tokens:
            if current_length + 1 > max_length - 2:  # Account for special tokens
                # Save current chunk
                if current_chunk:
                    chunks.append(" ".join(current_chunk))
                # Start new chunk
                current_chunk = [token]
                current_length = 1
            else:
                current_chunk.append(token)
                current_length += 1
        
        # Add the last chunk if not empty
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        
        logger.debug(f"Split text into {len(chunks)} chunks")
        return chunks
        
    except Exception as e:
        logger.warning(f"Error chunking text: {e}")
        return []

def detect_pii(text, scan_type="full"):
    """Detect PII in text using GLiNER model."""
    try:
        if not gliner_model:
            return []

        labels_to_use = PII_LABELS_FULL if scan_type == "full" else PII_LABELS
        
        # Use the new suppressor
        with SuppressStdoutStderr():
            entities = gliner_model.predict_entities(text, labels_to_use)
        
        formatted_entities = []
        for entity in entities:
            formatted_entities.append({
                "text": entity["text"],
                "label": entity["label"]
            })
            
        return formatted_entities
        
    except Exception as e:
        logger.warning(f"Error detecting PII: {e}")
        return []

def scan_file_for_pii(file_path, scan_type):
    """Scan a file for PII entities."""
    try:
        text = extract_text_from_file(file_path, scan_type)
        if not text:
            logger.error(f"Failed to extract text from {file_path}")
            return []

        chunks = chunk_text(text, MAX_CHUNK_LENGTH)
        all_pii_entities = []

        logger.info(f"Processing chunks for file: {file_path} ({scan_type})")
        for i, chunk in enumerate(chunks):
            logger.debug(f"Processing chunk {i + 1}/{len(chunks)}")
            pii_entities = detect_pii(chunk, scan_type)
            if pii_entities:
                all_pii_entities.extend(pii_entities)

        if all_pii_entities:
            # Get unique labels
            labels = ", ".join(sorted(set(entity['label'] for entity in all_pii_entities)))
            
            # Log detailed message
            log_message = f"PII data potentially exposed in {file_path} ({scan_type}):"
            for entity in all_pii_entities:
                log_message += f"\n  - {entity['label']}: {entity['text']}"
            logger.warning(log_message)
            
            # Print messages for Veeam detection
            print("PII data potentially exposed")  # For Veeam regex match
            print(f"PII_DETECTED: {labels}")  # For structured output
            
        return all_pii_entities
    except Exception as e:
        logger.error(f"Error scanning file {file_path}: {e}")
        sys.exit(EXIT_PII_DETECTION_ERROR)

def process_file(file_path, scan_type):
    """Process a single file for PII scanning."""
    logger.info(f"\nProcessing file: {file_path}")
    
    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension not in ALLOWED_FILE_TYPES:
        logger.warning(f"Skipping unsupported file type: {file_extension}")
        return

    file_size = os.path.getsize(file_path)
    file_modified = datetime.fromtimestamp(os.path.getmtime(file_path), tz=timezone.utc).isoformat()
    
    logger.info("Calculating checksum...")
    file_checksum = calculate_checksum(file_path, scan_type)

    if file_checksum is None:
        logger.warning(f"Skipping {file_path} due to checksum error.")
        return

    logger.info(f"Checksum: {file_checksum}")

    already_scanned, previous_pii_entities = is_file_scanned(file_path, file_checksum, scan_type)

    if already_scanned:
        if previous_pii_entities:
            labels = ", ".join(sorted(set(entity['label'] for entity in previous_pii_entities)))
            logger.warning(f"PII data potentially exposed in previously scanned file: {file_path} ({scan_type}): {labels}")
            print("PII data potentially exposed")
            print(f"PII_DETECTED: {labels}")
        else:
            logger.info(f"Skipping already scanned file: {file_path} ({scan_type}) - No PII found in previous scan")
        return

    logger.info(f"Scanning file for PII: {file_path} ({scan_type})")
    pii_entities = scan_file_for_pii(file_path, scan_type)

    logger.info("Saving results to database...")
    save_scan_result(file_path, pii_entities, file_size, file_modified, file_checksum, scan_type)

def scan_directory(directory, scan_type):
    """Scan all supported files in a directory."""
    for root, _, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            process_file(file_path, scan_type)

# Add a function to verify database
def verify_database():
    """Verify database exists and is properly initialized."""
    conn = None
    try:
        if not os.path.exists(DB_FILE):
            logger.debug(f"Database file not found at: {DB_FILE}")  # Changed to debug level
            init_db()
            return
            
        conn = sqlite3.connect(DB_FILE)
        cursor = conn.cursor()
        
        # Check if table exists
        cursor.execute("""
            SELECT name FROM sqlite_master 
            WHERE type='table' AND name='scan_history'
        """)
        
        if not cursor.fetchone():
            logger.debug("Database exists but missing required table.")  # Changed to debug level
            if conn:
                conn.close()
            init_db()
        else:
            logger.debug("Database verified successfully.")  # Changed to debug level
            
    except sqlite3.Error as e:
        logger.warning(f"Database verification error: {e}")  # Changed to warning level
        sys.exit(EXIT_DB_ERROR)
    finally:
        if conn:
            conn.close()

def custom_formatwarning(message, category, filename, lineno, line=None):
    """Custom format for UserWarning, logs it as INFO."""
    logger.info(f"{filename}:{lineno}: {category.__name__}: {message}")
    return ''  # Suppress the original warning

warnings.formatwarning = custom_formatwarning

if __name__ == "__main__":
    try:
        verify_database()
        init_nltk()
        
        parser = argparse.ArgumentParser(description="PII Scanner with Lite and Full Scan Options")
        parser.add_argument("path", help="The directory to scan")
        parser.add_argument("--scan-type", choices=["lite", "full"], default="full",
                          help="Specify 'lite' for a quick 1MB scan or 'full' for a complete scan (default: full)")
        parser.add_argument("--verbose", action="store_true", help="Enable verbose logging (DEBUG level)")

        args = parser.parse_args()
        
        if not args.path:
            logger.error("Missing file path argument")
            sys.exit(EXIT_MISSING_PATH)
            
        if args.scan_type not in ["lite", "full"]:
            logger.error("Invalid scan type specified")
            sys.exit(EXIT_INVALID_SCAN_TYPE)

        if args.verbose:
            logger.setLevel(logging.DEBUG)
            logger.debug("Verbose logging enabled")

        # Initialize models
        if gliner_model is None:
            logger.error("Failed to initialize required models")
            sys.exit(EXIT_GLINER_INIT_ERROR)

        # Scan directory using process_file
        pii_found = False
        for root, _, files in os.walk(args.path):
            for file in files:
                file_path = os.path.join(root, file)
                try:
                    if os.path.splitext(file)[1].lower() not in ALLOWED_FILE_TYPES:
                        continue
                        
                    # Use process_file instead of scan_file_for_pii directly
                    process_file(file_path, args.scan_type)
                    
                    # Check if PII was found by querying the database
                    conn = sqlite3.connect(DB_FILE)
                    cursor = conn.cursor()
                    cursor.execute("""
                        SELECT pii_entities 
                        FROM scan_history 
                        WHERE file_path = ? AND scan_type = ?
                    """, (file_path, args.scan_type))
                    
                    result = cursor.fetchone()
                    if result and result[0] != '[]':
                        pii_found = True
                    
                    conn.close()
                        
                except FileNotFoundError:
                    logger.error(f"File not found: {file_path}")
                    sys.exit(EXIT_FILE_NOT_FOUND)
                except Exception as e:
                    logger.error(f"Error processing file {file_path}: {e}")
                    sys.exit(EXIT_GENERAL_ERROR)

        logger.info("Scanning complete.")
        sys.exit(EXIT_PII_FOUND if pii_found else EXIT_SUCCESS)

    except Exception as e:
        logger.error(f"An error occurred: {e}")
        sys.exit(EXIT_GENERAL_ERROR) 